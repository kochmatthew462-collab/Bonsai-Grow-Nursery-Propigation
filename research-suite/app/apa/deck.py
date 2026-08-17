"""
APA-compliant PowerPoint decks.

APA 7 has no chapter on slides, so what "APA-compliant slides" means in practice
is the set of rules a course rubric or a defence committee actually checks:

* **In-text citations appear on the slide itself**, next to the borrowed content,
  not gathered on a references slide at the end. A slide with an uncited
  statistic is the presentation equivalent of an uncited sentence.
* **A References slide (or several) at the end**, entries in full APA form with
  hanging indents, in the same alphabetical order as the paper's reference list.
* **Figures carry their number, italic title and note**, including source
  attribution, exactly as in the paper.
* **Text is large enough to read from the back of a room** — 28 pt body minimum
  here, 40 pt titles, which is well above the 18 pt that looks fine on a laptop
  and is illegible projected.

The deck is built from the same claim ledger as the paper, so a slide's citation
cannot disagree with the manuscript's — they are generated from one source of
truth by one citation engine.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..models import Project, SupportType, Work
from .citations import CitationContext, Run, intext, reference_list

# 16:9, the shape every projector and conference room expects.
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN = Inches(0.75)
TITLE_SIZE = Pt(40)
BODY_SIZE = Pt(28)
SMALL_SIZE = Pt(18)
REFERENCE_SIZE = Pt(14)

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)

# A references slide holds this many entries before spilling to the next. Set
# from what actually fits at 14 pt with a hanging indent, not from a guess.
REFERENCES_PER_SLIDE = 5
BULLETS_PER_SLIDE = 5


@dataclass
class Slide:
    """One planned slide."""

    title: str
    bullets: list[tuple[str, list[Run]]] = field(default_factory=list)
    figure_path: str = ""
    figure_number: int = 0
    figure_title: str = ""
    figure_note: str = ""
    speaker_notes: str = ""
    kind: str = "content"


class ApaDeck:
    """Builds an APA-conventional deck with python-pptx."""

    def __init__(self, project: Project, context: CitationContext,
                 font: str = "Calibri"):
        self.project = project
        self.context = context
        self.font = font
        self.presentation = Presentation()
        self.presentation.slide_width = SLIDE_WIDTH
        self.presentation.slide_height = SLIDE_HEIGHT
        self._figure_number = 0

    # ------------------------------------------------------------ primitives

    def _blank(self):
        # Layout 6 is the blank layout — placeholders from the themed layouts
        # bring their own fonts and sizes, which then fight anything set here.
        return self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    def _textbox(self, slide, left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        return frame

    def _write(self, paragraph, runs: list[Run], size, *, bold=False,
               colour=INK, italic_override=None):
        for spec in runs:
            if not spec.text:
                continue
            run = paragraph.add_run()
            run.text = spec.text
            run.font.size = size
            run.font.bold = bold
            run.font.italic = spec.italic if italic_override is None else italic_override
            run.font.name = self.font
            run.font.color.rgb = colour

    def _slide_title(self, slide, text: str) -> None:
        frame = self._textbox(slide, MARGIN, Inches(0.45),
                              SLIDE_WIDTH - 2 * MARGIN, Inches(1.0))
        paragraph = frame.paragraphs[0]
        self._write(paragraph, [Run(text)], TITLE_SIZE, bold=True)

    def _notes(self, slide, text: str) -> None:
        if text:
            slide.notes_slide.notes_text_frame.text = text

    # ------------------------------------------------------------ title slide

    def title_slide(self) -> None:
        page = self.project.title_page
        slide = self._blank()
        frame = self._textbox(slide, MARGIN, Inches(2.1),
                              SLIDE_WIDTH - 2 * MARGIN, Inches(3.2))
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        self._write(paragraph, [Run(page.title or self.project.topic)],
                    Pt(44), bold=True)

        for line in (
            list(page.authors) + list(page.affiliations)
            + ([page.course, page.instructor, page.due_date]
               if page.variant == "student" else [])
        ):
            if not line:
                continue
            extra = frame.add_paragraph()
            extra.alignment = PP_ALIGN.CENTER
            extra.space_before = Pt(10)
            self._write(extra, [Run(line)], BODY_SIZE, colour=INK_SECONDARY)

        self._notes(slide, "Title slide. APA does not specify slide layouts; this "
                           "carries the same identifying details as the paper's "
                           "title page.")

    # --------------------------------------------------------- content slides

    def content_slide(self, plan: Slide) -> None:
        slide = self._blank()
        self._slide_title(slide, plan.title)

        frame = self._textbox(slide, MARGIN, Inches(1.75),
                              SLIDE_WIDTH - 2 * MARGIN, Inches(4.8))
        first = True
        for _, runs in plan.bullets:
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            paragraph.space_after = Pt(16)
            paragraph.level = 0
            # A literal bullet glyph: python-pptx cannot set list formatting on a
            # plain textbox, and a themed placeholder would override the fonts.
            self._write(paragraph, [Run("•  ")], BODY_SIZE, colour=ACCENT)
            self._write(paragraph, runs, BODY_SIZE)

        self._notes(slide, plan.speaker_notes)

    def figure_slide(self, plan: Slide) -> None:
        slide = self._blank()
        self._slide_title(slide, plan.title)

        self._figure_number += 1
        number = plan.figure_number or self._figure_number

        caption = self._textbox(slide, MARGIN, Inches(1.6),
                                SLIDE_WIDTH - 2 * MARGIN, Inches(0.8))
        paragraph = caption.paragraphs[0]
        self._write(paragraph, [Run(f"Figure {number}")], SMALL_SIZE, bold=True)
        if plan.figure_title:
            second = caption.add_paragraph()
            self._write(second, [Run(plan.figure_title, italic=True)], SMALL_SIZE)

        if plan.figure_path and Path(plan.figure_path).exists():
            available_height = Inches(3.5)
            picture = slide.shapes.add_picture(
                plan.figure_path, MARGIN, Inches(2.5), height=available_height)
            # Centre horizontally, and shrink to fit if the aspect ratio makes it
            # wider than the slide.
            max_width = SLIDE_WIDTH - 2 * MARGIN
            if picture.width > max_width:
                ratio = max_width / picture.width
                picture.width = Emu(int(picture.width * ratio))
                picture.height = Emu(int(picture.height * ratio))
            picture.left = Emu(int((SLIDE_WIDTH - picture.width) / 2))

        if plan.figure_note:
            note = self._textbox(slide, MARGIN, Inches(6.25),
                                 SLIDE_WIDTH - 2 * MARGIN, Inches(0.9))
            paragraph = note.paragraphs[0]
            self._write(paragraph, [Run("Note. ")], Pt(12), italic_override=True,
                        colour=INK_SECONDARY)
            self._write(paragraph, [Run(plan.figure_note)], Pt(12),
                        colour=INK_SECONDARY)

        self._notes(slide, plan.speaker_notes)

    # ------------------------------------------------------ references slides

    def reference_slides(self, works: list[Work]) -> None:
        """References in full APA form, spilling across as many slides as needed."""
        if not works:
            return
        entries = reference_list(works, self.context)
        for start in range(0, len(entries), REFERENCES_PER_SLIDE):
            chunk = entries[start:start + REFERENCES_PER_SLIDE]
            slide = self._blank()
            continued = " (continued)" if start else ""
            self._slide_title(slide, f"References{continued}")

            frame = self._textbox(slide, MARGIN, Inches(1.7),
                                  SLIDE_WIDTH - 2 * MARGIN, Inches(5.2))
            for index, entry in enumerate(chunk):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.space_after = Pt(14)
                _hanging_indent(paragraph, Inches(0.5))
                self._write(paragraph, entry, REFERENCE_SIZE)
            self._notes(slide, "Reference entries are identical to the paper's "
                               "reference list and in the same order.")

    # ------------------------------------------------------------------- plan

    def plan_from_claims(self, *, max_bullets: int = BULLETS_PER_SLIDE) -> list[Slide]:
        """Build a slide plan from the project's claim ledger.

        One section becomes one or more slides; each claim becomes a bullet with
        its citation attached. Because the citations come from the same engine as
        the paper's, a slide can never cite something the paper does not.
        """
        self.context.reset_group_state()
        by_key = {w.key: w for w in self.project.works}
        slides: list[Slide] = []

        sections: list[tuple[str, list]] = []
        for claim in sorted(self.project.claims, key=lambda c: c.order):
            if sections and sections[-1][0] == claim.section:
                sections[-1][1].append(claim)
            else:
                sections.append((claim.section, [claim]))

        for heading, claims in sections:
            chunks = [claims[i:i + max_bullets]
                      for i in range(0, len(claims), max_bullets)]
            for index, chunk in enumerate(chunks):
                title = heading + (" (continued)" if index else "")
                plan = Slide(title=title)
                for claim in chunk:
                    works = [by_key[k] for k in claim.work_keys if k in by_key]
                    runs = [Run(_to_bullet(claim.text))]
                    if works and claim.support_type is not SupportType.NO_CITATION:
                        runs.append(Run(" "))
                        runs.extend(intext(works, self.context,
                                           locator=claim.locus or ""))
                    plan.bullets.append((claim.claim_id, runs))
                plan.speaker_notes = "\n".join(
                    f"{c.claim_id}: {c.rationale}" for c in chunk if c.rationale
                )
                slides.append(plan)
        return slides

    def build(self, slides: list[Slide] | None = None,
              figures: list[Slide] | None = None) -> "ApaDeck":
        self.title_slide()
        for plan in (slides if slides is not None else self.plan_from_claims()):
            self.content_slide(plan)
        for plan in figures or []:
            self.figure_slide(plan)
        self.reference_slides(self.project.cited_works())
        return self

    def save(self, path: str) -> str:
        self.presentation.save(path)
        return path


def _hanging_indent(paragraph, size: Emu) -> None:
    """Apply a hanging indent to a pptx paragraph.

    Unlike python-docx, python-pptx exposes no `paragraph_format` — indentation
    lives in the DrawingML `a:pPr` element as `marL` (left margin) and `indent`
    (first-line offset, negative for a hanging indent), both in EMU. Setting them
    directly is the only route.
    """
    properties = paragraph._p.get_or_add_pPr()
    properties.set("marL", str(int(size)))
    properties.set("indent", str(-int(size)))


def _to_bullet(text: str) -> str:
    """Shorten a paper sentence into slide text.

    Slides are not paragraphs. The assembler markers are stripped and the
    trailing period dropped — a bulleted fragment does not take one — but the
    wording is otherwise the author's, because rewriting it here would put text
    on the slide that no claim in the ledger accounts for.
    """
    cleaned = text.replace("[[CITE]]", "").replace("[[NARRATIVE]]", "").strip()
    return cleaned.rstrip(".").strip()


def build_deck(
    project: Project, context: CitationContext, path: str,
    figures: list[Slide] | None = None, font: str = "Calibri",
) -> str:
    return ApaDeck(project, context, font).build(figures=figures).save(path)
